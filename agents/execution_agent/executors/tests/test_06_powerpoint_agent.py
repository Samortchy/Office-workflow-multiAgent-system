"""
Agent 06 — PowerPoint Agent
approval: single_confirm  |  on_failure: escalate

Known config issues (do NOT fix here — flag to agent-owner):
  - DBFetcher step uses unsupported 'match_on' key (DBFetcher uses 'filters') → silently ignored.
  - DBFetcher table "templates/pptx/templates_meta.json" is not a valid SQLite table name.
    DBFetcher falls back to empty rows (DB absent) without error — step succeeds.
  - prompt_template "generate_slides" is not in LLMGenerator._PROMPTS; the raw key
    string is used as the prompt.  LLMGenerator._call is mocked regardless.

test_06_powerpoint_agent / test_06_powerpoint_agent_confirmed:
  PPTXWriter is mocked (CI-safe, no python-pptx dependency needed).

test_06_powerpoint_agent_real_file:
  PPTXWriter runs for real — produces an actual .pptx file on disk.
  Only LLMGenerator._call is mocked (no API key required).
  The file is verified with python-pptx then deleted after the test.
"""
import json
from unittest.mock import patch
from core.base_agent import ExecutionRunner
from steps.base_step import StepResult

_ENVELOPE = {
    "intake": {
        "department": "Finance",
        "task_type": "presentation",
        "isAutonomous": False,
        "confidence": 0.85,
        "processed_at": "2026-05-03T10:00:00Z",
    },
    "task": {
        "task_id": "T006",
        "title": "Q2 Board Review Presentation",
        "description": "Create a PowerPoint for the Q2 board review meeting.",
        "department": "Finance",
        "isAutonomous": False,
        "task_type": "presentation",
        "requester_name": "eve@company.com",
        "stated_deadline": "2026-05-07",
        "action_required": "Generate and email the presentation file",
        "success_criteria": "Presentation sent to requester",
        "structured_at": "2026-05-03T10:05:00Z",
    },
    "priority": {
        "priority_score": 3,
        "priority_label": "medium",
        "confidence": 0.85,
        "model_version": "v1",
        "top_features_used": [],
        "scored_at": "2026-05-03T10:10:00Z",
    },
}

_SLIDE_JSON = json.dumps({
    "slides": [
        {
            "title": "Q2 Financial Overview",
            "bullet_points": ["Revenue: 12.4M EGP", "vs Q1: +8%"],
            "speaker_notes": "Highlight the Q1-to-Q2 improvement.",
        }
    ],
    "template_path": "",
    "presentation_title": "Q2 Board Review",
    "paused_for_clarification": False,
})


_PPTX_RESULT = StepResult(
    True,
    {
        "output_path": "output/presentations/T006_20260503_q2_board_review.pptx",
        "slides_written": 1,
        "template_used": "blank",
        "paused": False,
        "clarification_question": "",
    },
    None,
)


def test_06_powerpoint_agent():
    runner = ExecutionRunner("configs/06_powerpoint_agent.json")

    with (
        patch(
            "steps.processors.llm_generator.LLMGenerator._call",
            return_value=_SLIDE_JSON,
        ),
        patch("steps.custom.pptx_writer.PPTXWriter.run", return_value=_PPTX_RESULT),
    ):
        result = runner.execute(_ENVELOPE.copy())

    assert "execution" in result
    assert result["execution"]["agent_name"] == "powerpoint_agent"
    # approval=single_confirm → runner pauses before EmailDispatcher
    assert result["execution"]["status"] == "approval_pending"
    assert "generate_slide_json" in result["execution"]["steps"]
    assert "write_pptx" in result["execution"]["steps"]


def test_06_powerpoint_agent_confirmed():
    """
    Mimics user reviewing and approving the generated presentation.

    Phase 1 → pauses (single_confirm gate fires before EmailDispatcher).
    Phase 2 → re-run with the paused envelope → gate skipped → EmailDispatcher
    sends the .pptx as a dry-run email attachment.

    Both runs call LLMGenerator._call and PPTXWriter.run (same mocks active).
    """
    runner = ExecutionRunner("configs/06_powerpoint_agent.json")

    with (
        patch(
            "steps.processors.llm_generator.LLMGenerator._call",
            return_value=_SLIDE_JSON,
        ),
        patch("steps.custom.pptx_writer.PPTXWriter.run", return_value=_PPTX_RESULT),
    ):
        # Phase 1 — pauses for approval
        paused = runner.execute(_ENVELOPE.copy())
        assert paused["execution"]["status"] == "approval_pending"

        # Phase 2 — mimic confirmation
        result = runner.execute(paused)

    assert result["execution"]["agent_name"] == "powerpoint_agent"
    assert result["execution"]["status"] == "completed"
    assert "email_file" in result["execution"]["steps"]


# ---------------------------------------------------------------------------
# Rich slide JSON used by the real-file test (3 slides, proper content)
# ---------------------------------------------------------------------------
_REAL_SLIDE_JSON = json.dumps({
    "slides": [
        {
            "title": "Q2 Financial Overview",
            "bullet_points": [
                "Total revenue: 12.4M EGP",
                "Quarter-on-quarter growth: +8%",
                "Budget variance: -2%",
            ],
            "speaker_notes": "Open with the headline number, then pivot to growth story.",
        },
        {
            "title": "Expense Breakdown",
            "bullet_points": [
                "Salaries & benefits: 65%",
                "Operations & logistics: 20%",
                "Miscellaneous: 15%",
            ],
            "speaker_notes": "Highlight the ops efficiency gain vs Q1.",
        },
        {
            "title": "Q3 Outlook & Next Steps",
            "bullet_points": [
                "Revenue target: 13.5M EGP (+9%)",
                "Cost reduction programme underway",
                "Two new client contracts in pipeline",
            ],
            "speaker_notes": "End on a forward-looking note; invite questions.",
        },
    ],
    "template_path": "",
    "presentation_title": "Q2 Board Review",
    "paused_for_clarification": False,
})


def test_06_powerpoint_agent_real_file():
    """
    End-to-end test: PPTXWriter runs for real and produces an actual .pptx file.

    Flow
    ----
    Phase 1 — runner pauses at the approval gate (single_confirm).
    Phase 2 — confirmation re-run: PPTXWriter writes a real PPTX to
               output/presentations/, EmailDispatcher attaches it in dry-run.

    Only LLMGenerator._call is mocked (avoids the OpenRouter API call while
    letting the full PPTXWriter + EmailDispatcher pipeline execute).

    Assertions
    ----------
    - Pipeline reaches status="completed"
    - The output_path returned by PPTXWriter points to a real file on disk
    - The file is non-empty and openable as a valid PPTX
    - The presentation contains exactly 3 slides with the expected titles
    - The file is removed after the test (teardown via finally)
    """
    from pathlib import Path
    from pptx import Presentation

    runner = ExecutionRunner("configs/06_powerpoint_agent.json")
    output_file = None

    try:
        with patch(
            "steps.processors.llm_generator.LLMGenerator._call",
            return_value=_REAL_SLIDE_JSON,
        ):
            # Phase 1 — pauses for user approval
            paused = runner.execute(_ENVELOPE.copy())
            assert paused["execution"]["status"] == "approval_pending"

            # Phase 2 — mimic confirmation; PPTXWriter now runs for real
            result = runner.execute(paused)

        # ── Pipeline assertions ──────────────────────────────────────────
        assert result["execution"]["agent_name"] == "powerpoint_agent"
        assert result["execution"]["status"] == "completed"
        assert "write_pptx" in result["execution"]["steps"]
        assert "email_file" in result["execution"]["steps"]

        pptx_data = result["execution"]["steps"]["write_pptx"]["data"]
        assert pptx_data["paused"] is False
        assert pptx_data["slides_written"] == 3

        # ── File-system assertions ───────────────────────────────────────
        output_file = Path(pptx_data["output_path"])
        assert output_file.exists(), f"PPTX file not found at {output_file}"
        assert output_file.stat().st_size > 0, "PPTX file is empty"

        # ── Content assertions (open with python-pptx) ──────────────────
        prs = Presentation(str(output_file))
        assert len(prs.slides) == 3, f"Expected 3 slides, got {len(prs.slides)}"

        titles = [
            slide.shapes.title.text
            for slide in prs.slides
            if slide.shapes.title
        ]
        assert titles[0] == "Q2 Financial Overview"
        assert titles[1] == "Expense Breakdown"
        assert titles[2] == "Q3 Outlook & Next Steps"

        # ── Email dry-run assertion ──────────────────────────────────────
        email_data = result["execution"]["steps"]["email_file"]["data"]
        assert email_data["dry_run"] is True
        assert str(output_file) in email_data["output_path"] or output_file.exists()

    finally:
        pass  # file is intentionally kept for inspection
