"""
pptx_writer.py
Custom step — PowerPoint Agent (Agent 06)
P5 responsibility

What it does:
    Receives a structured slide JSON produced by the preceding LLMGenerator step
    and writes it to a .pptx file using python-pptx and a pre-existing template.

    The LLMGenerator step before this one (named "generate_slide_json") produces:
        {
            "template_path": "templates/pptx/finance_board_review.pptx",
            "slides": [
                {
                    "title":         "Q2 Financial Overview",
                    "bullet_points": ["Revenue: 12.4M EGP", "vs Q1: +8%", "vs Budget: -2%"],
                    "speaker_notes": "Highlight the Q1-to-Q2 improvement story here."
                },
                ...
            ],
            "presentation_title": "Q2 Budget Review",
            "paused_for_clarification": false
        }

    If "paused_for_clarification" is True in the LLM output, this step skips
    writing and returns a paused signal so the runner can pause execution.

Output file naming (from spec config):
    {task_id}_{date}_{short_title}.pptx
    Example: TASK-2441_20260430_Q2_budget_review.pptx

Returns:
    StepResult.data = {
        "output_path":    str,          # full path to the written .pptx file
        "slides_written": int,
        "template_used":  str,
        "paused":         bool,         # True if LLM asked for clarification
        "clarification_question": str,  # populated if paused == True, else ""
    }

    If paused == True, the runner reads data["status"] = "paused" and halts
    execution to await user input. No .pptx file is written in this case.

Spec compliance:
    - Inherits BaseStep
    - run() signature: (self, envelope: dict, config: dict) -> StepResult
    - Never raises — all exceptions caught and returned as StepResult(success=False)
    - Never modifies envelope directly
    - Never adds fields to StepResult
    - Never writes to SQLite — runner does that
    - Uses resolve_path() from core.envelope for all envelope reads
    - Writes files only to the sandboxed output_dir from config — never arbitrary paths
"""

import os
import re
from datetime import date
from steps.base_step import BaseStep, StepResult
from core.envelope import resolve_path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PPTXWriter(BaseStep):

    def run(self, envelope: dict, config: dict) -> StepResult:
        try:
            # ── 0. Dependency check ───────────────────────────────────────
            if not PPTX_AVAILABLE:
                return StepResult(
                    success=False,
                    data={},
                    error=(
                        "PPTXWriter: python-pptx is not installed. "
                        "Run: pip install python-pptx"
                    )
                )

            # ── 1. Read config ────────────────────────────────────────────
            output_dir    = config.get("output_dir", "output/presentations")
            naming_pattern= config.get("naming", "{task_id}_{date}_{short_title}.pptx")

            # ── 2. Read slide JSON from the LLM generator step ────────────
            slide_data = resolve_path(
                envelope,
                "execution.steps.generate_slide_json.data"
            )

            # ── 3. Check if LLM is waiting for clarification ──────────────
            if slide_data.get("paused_for_clarification", False):
                return StepResult(
                    success=True,
                    data={
                        "output_path":             "",
                        "slides_written":          0,
                        "template_used":           "",
                        "paused":                  True,
                        "clarification_question":  slide_data.get(
                            "clarification_question", "Could you provide more details?"
                        ),
                        "status": "paused",   # runner reads this to pause execution
                    },
                    error=None
                )

            slides            = slide_data.get("slides", [])
            template_path     = slide_data.get("template_path", "")
            presentation_title= slide_data.get("presentation_title", "Presentation")

            if not slides:
                return StepResult(
                    success=False,
                    data={},
                    error="PPTXWriter: no slides found in generate_slide_json output."
                )

            # ── 4. Validate template path is within templates/pptx/ ───────
            # Security: never write to arbitrary paths from user input
            if template_path and not template_path.startswith("templates/pptx/"):
                return StepResult(
                    success=False,
                    data={},
                    error=(
                        f"PPTXWriter: template_path '{template_path}' is outside "
                        f"the allowed directory 'templates/pptx/'. Refusing to load."
                    )
                )

            # ── 5. Load template or create blank presentation ─────────────
            if template_path and os.path.exists(template_path):
                prs = Presentation(template_path)
            else:
                prs = Presentation()

            # ── 6. Write slides ───────────────────────────────────────────
            slides_written = 0
            for slide_spec in slides:
                self._write_slide(prs, slide_spec)
                slides_written += 1

            # ── 7. Build output filename ──────────────────────────────────
            task_id     = envelope.get("task", {}).get("task_id", "TASK-UNKNOWN")
            today       = date.today().strftime("%Y%m%d")
            short_title = self._slugify(presentation_title)[:30]

            filename = (
                naming_pattern
                .replace("{task_id}",     task_id)
                .replace("{date}",        today)
                .replace("{short_title}", short_title)
            )

            # ── 8. Validate output_dir is within allowed sandbox ──────────
            if not output_dir.startswith("output/"):
                return StepResult(
                    success=False,
                    data={},
                    error=(
                        f"PPTXWriter: output_dir '{output_dir}' is outside "
                        f"the allowed sandbox 'output/'. Refusing to write."
                    )
                )

            os.makedirs(output_dir, exist_ok=True)
            output_path = os.path.join(output_dir, filename)
            prs.save(output_path)

            return StepResult(
                success=True,
                data={
                    "output_path":            output_path,
                    "slides_written":         slides_written,
                    "template_used":          template_path or "blank",
                    "paused":                 False,
                    "clarification_question": "",
                },
                error=None
            )

        except KeyError as e:
            return StepResult(
                success=False,
                data={},
                error=(
                    f"PPTXWriter could not find required envelope path: {e}. "
                    f"Ensure 'generate_slide_json' step ran successfully before this step."
                )
            )
        except Exception as e:
            return StepResult(
                success=False,
                data={},
                error=f"PPTXWriter unexpected error: {str(e)}"
            )

    # ── private helpers ───────────────────────────────────────────────────────

    def _write_slide(self, prs: "Presentation", slide_spec: dict) -> None:
        """
        Adds one slide to the presentation.

        Uses slide layout index 1 (Title and Content) if available,
        falls back to layout 0 (Title Slide) if the template has fewer layouts.

        slide_spec shape:
            {
                "title":         str,
                "bullet_points": list[str],
                "speaker_notes": str   (optional)
            }
        """
        title_text    = slide_spec.get("title", "")
        bullet_points = slide_spec.get("bullet_points", [])
        speaker_notes = slide_spec.get("speaker_notes", "")

        # Pick the best available layout
        layout_index = 1 if len(prs.slide_layouts) > 1 else 0
        slide_layout = prs.slide_layouts[layout_index]
        slide        = prs.slides.add_slide(slide_layout)

        # Write title
        if slide.shapes.title:
            slide.shapes.title.text = title_text

        # Write body / bullet points
        # Find the first non-title placeholder to use as the content area
        body_placeholder = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:   # idx 0 is always the title
                body_placeholder = ph
                break

        if body_placeholder and bullet_points:
            tf = body_placeholder.text_frame
            tf.clear()
            for i, point in enumerate(bullet_points):
                if i == 0:
                    tf.paragraphs[0].text = point
                else:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0

        # Write speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def _slugify(self, text: str) -> str:
        """
        Converts a presentation title to a safe filename fragment.
        Example: "Q2 Budget Review — Board" → "Q2_budget_review_board"
        """
        text = text.lower()
        text = re.sub(r"[^a-z0-9\s]", "", text)   # remove non-alphanumeric
        text = re.sub(r"\s+", "_", text.strip())   # spaces to underscores
        return text